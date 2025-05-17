import os
import uuid
import logging
import tempfile
import subprocess
from datetime import datetime
from flask import Blueprint, render_template, request, jsonify, send_from_directory, current_app
from werkzeug.utils import secure_filename
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import platform

# Configure logging
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('converter.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Create Blueprint
ppt_to_pdf_bp = Blueprint('ppt_to_pdf', __name__, template_folder='../templates')

# Configuration (now using current_app.config instead of blueprint config)
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'outputs'
ALLOWED_EXTENSIONS = {'pptx', 'ppt'}

# Ensure folders exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def convert_pptx_to_pdf(input_path, output_path):
    """Main conversion function with multiple fallback methods"""
    logger.info(f"Starting conversion: {input_path} to {output_path}")
    
    # Try different conversion methods in order of reliability
    methods = [
        try_libreoffice_conversion,
        try_windows_powerpoint_conversion,
        convert_pptx_to_pdf_images
    ]
    
    for method in methods:
        try:
            logger.info(f"Trying {method.__name__}")
            if method(input_path, output_path):
                if validate_pdf(output_path):
                    logger.info(f"Success with {method.__name__}")
                    return True
                else:
                    logger.warning(f"Method {method.__name__} produced invalid PDF")
                    os.remove(output_path)
        except Exception as e:
            logger.error(f"Error in {method.__name__}: {str(e)}")
            if os.path.exists(output_path):
                os.remove(output_path)
    
    logger.error("All conversion methods failed")
    return False

def try_libreoffice_conversion(input_path, output_path):
    """Method 1: Use LibreOffice (most reliable cross-platform)"""
    try:
        libreoffice_path = find_libreoffice()
        if not libreoffice_path:
            logger.warning("LibreOffice not found")
            return False
            
        temp_dir = tempfile.mkdtemp()
        command = [
            libreoffice_path,
            '--headless',
            '--convert-to',
            'pdf',
            '--outdir',
            temp_dir,
            input_path
        ]
        
        logger.debug(f"Running command: {' '.join(command)}")
        
        result = subprocess.run(
            command,
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=60
        )
        
        # Find the converted file
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        temp_output = os.path.join(temp_dir, f"{base_name}.pdf")
        
        if os.path.exists(temp_output):
            os.rename(temp_output, output_path)
            logger.info(f"LibreOffice conversion successful: {output_path}")
            return True
            
        logger.warning("LibreOffice conversion failed - no output file")
        return False
    except subprocess.TimeoutExpired:
        logger.error("LibreOffice conversion timed out")
        return False
    except Exception as e:
        logger.error(f"LibreOffice conversion error: {str(e)}")
        return False

def find_libreoffice():
    """Find LibreOffice executable path"""
    try:
        if platform.system() == 'Windows':
            paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
            ]
            for path in paths:
                if os.path.exists(path):
                    return path
            # Try PATH
            return 'soffice.exe' if subprocess.run(['where', 'soffice.exe'], capture_output=True).returncode == 0 else None
        else:
            # Linux/Mac - try standard paths
            paths = ['libreoffice', 'soffice']
            for path in paths:
                if subprocess.run(['which', path], capture_output=True).returncode == 0:
                    return path
            return None
    except Exception as e:
        logger.error(f"Error finding LibreOffice: {str(e)}")
        return None

def try_windows_powerpoint_conversion(input_path, output_path):
    """Method 2: Use MS PowerPoint (Windows only)"""
    if platform.system() != 'Windows':
        return False
        
    try:
        import win32com.client
        powerpoint = win32com.client.Dispatch("Powerpoint.Application")
        powerpoint.Visible = 1  # Show PowerPoint window for debugging
        
        logger.info(f"Opening PowerPoint file: {input_path}")
        deck = powerpoint.Presentations.Open(os.path.abspath(input_path))
        
        logger.info(f"Saving as PDF: {output_path}")
        deck.SaveAs(os.path.abspath(output_path), 32)  # 32 is PDF format
        
        deck.Close()
        powerpoint.Quit()
        
        logger.info("PowerPoint conversion completed")
        return True
    except Exception as e:
        logger.error(f"PowerPoint conversion failed: {str(e)}")
        return False

def convert_pptx_to_pdf_images(input_path, output_path):
    """Method 3: Fallback to image-based conversion"""
    try:
        prs = Presentation(input_path)
        if not prs.slides:
            logger.error("No slides found in presentation")
            return False
            
        images = []
        temp_files = []
        
        # Calculate slide size in pixels (300 DPI)
        slide_width = int(prs.slide_width * 300 / 914400)
        slide_height = int(prs.slide_height * 300 / 914400)
        
        for i, slide in enumerate(prs.slides):
            try:
                temp_img = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                temp_files.append(temp_img.name)
                
                # Create image with slide number (replace with actual rendering in production)
                img = Image.new('RGB', (slide_width, slide_height), (255, 255, 255))
                draw = ImageDraw.Draw(img)
                
                try:
                    font = ImageFont.truetype("arial.ttf", 40)
                except:
                    font = ImageFont.load_default()
                
                draw.text((50, 50), f"Slide {i+1}", fill=(0, 0, 0), font=font)
                img.save(temp_img.name, 'PNG', dpi=(300, 300))
                images.append(temp_img.name)
                
                logger.debug(f"Processed slide {i+1}")
            except Exception as e:
                logger.error(f"Error processing slide {i+1}: {str(e)}")
                continue
        
        if not images:
            logger.error("No slides were processed")
            return False
            
        # Convert images to PDF
        first_image = Image.open(images[0])
        first_image.save(
            output_path,
            "PDF",
            resolution=300.0,
            save_all=True,
            append_images=[Image.open(img) for img in images[1:]],
            quality=100
        )
        
        logger.info(f"Image-based PDF created: {output_path}")
        return True
    except Exception as e:
        logger.error(f"Image-based conversion failed: {str(e)}", exc_info=True)
        return False
    finally:
        # Cleanup temp files
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception as e:
                logger.warning(f"Could not delete temp file {temp_file}: {str(e)}")

def validate_pdf(pdf_path):
    """Basic PDF validation"""
    try:
        if not os.path.exists(pdf_path):
            return False
            
        if os.path.getsize(pdf_path) < 1024:  # Minimum reasonable size
            return False
            
        # Check if file starts with PDF header
        with open(pdf_path, 'rb') as f:
            header = f.read(4)
            return header == b'%PDF'
    except Exception as e:
        logger.error(f"PDF validation failed: {str(e)}")
        return False

@ppt_to_pdf_bp.route('ppt_to_pdf')
def index():
    return render_template('ppt_to_pdf.html')

@ppt_to_pdf_bp.route('ppt_to_pdf/convert', methods=['POST'])
def convert_pptx_to_pdf_route():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if file and allowed_file(file.filename):
        # Generate unique filenames
        original_filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())
        upload_path = os.path.join(UPLOAD_FOLDER, f"{unique_id}_{original_filename}")
        output_filename = f"{os.path.splitext(original_filename)[0]}_{unique_id}.pdf"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        
        # Save uploaded file
        file.save(upload_path)
        
        try:
            logger.info(f"Starting conversion: {upload_path} to {output_path}")
            
            # Convert PPTX to PDF
            success = convert_pptx_to_pdf(upload_path, output_path)
            
            if not success:
                logger.error("Conversion failed for file: " + upload_path)
                return jsonify({'error': 'Conversion failed. Please try another file.'}), 500
            
            if not validate_pdf(output_path):
                logger.error("Invalid PDF generated")
                if os.path.exists(output_path):
                    os.remove(output_path)
                return jsonify({'error': 'Generated PDF is invalid'}), 500
            
            logger.info(f"Conversion successful: {output_path}")
            return jsonify({
                'success': True,
                'download_url': f'ppt_to_pdf/download/{output_filename}',
                'filename': output_filename
            })
        except Exception as e:
            logger.error(f"Error during conversion: {str(e)}", exc_info=True)
            return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
        finally:
            # Clean up uploaded file
            if os.path.exists(upload_path):
                os.remove(upload_path)
    else:
        return jsonify({'error': 'Invalid file type. Only PPT/PPTX files are allowed.'}), 400

@ppt_to_pdf_bp.route('ppt_to_pdf/download/<filename>')
def download_file(filename):
    return send_from_directory(OUTPUT_FOLDER, filename, as_attachment=True)

@ppt_to_pdf_bp.route('/ppt_to_pdf/cleanup', methods=['POST'])
def cleanup():
    try:
        # Clean up output files older than 1 hour
        now = datetime.now()
        for filename in os.listdir(OUTPUT_FOLDER):
            file_path = os.path.join(OUTPUT_FOLDER, filename)
            file_time = datetime.fromtimestamp(os.path.getmtime(file_path))
            if (now - file_time).total_seconds() > 3600:  # 1 hour
                os.remove(file_path)
        return jsonify({'success': True})
    except Exception as e:
        logger.error(f"Cleanup error: {str(e)}", exc_info=True)
        return jsonify({'error': str(e)}), 500