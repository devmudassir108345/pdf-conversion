# Add this at the top of app.py
import logging
logging.basicConfig(level=logging.DEBUG)
import os
from flask import Blueprint, render_template, request, jsonify, send_from_directory
from werkzeug.utils import secure_filename
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import uuid
from datetime import datetime
from PIL import Image
import tempfile

pdf_to_ppt_bp = Blueprint('pdf_to_ppt', __name__, template_folder='../templates')

# Configuration







pdf_to_ppt_bp.upload_folder = 'uploads'
pdf_to_ppt_bp.output_folder = 'outputs'
pdf_to_ppt_bp.allowed_extensions = {'pdf'}

# Ensure folders exist
os.makedirs(pdf_to_ppt_bp.upload_folder, exist_ok=True)
os.makedirs(pdf_to_ppt_bp.output_folder, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in pdf_to_ppt_bp.allowed_extensions





# Alternative conversion function using PyMuPDF
def convert_pdf_to_pptx(pdf_path, output_path):
    try:
        import fitz  # PyMuPDF
        from io import BytesIO
        
        doc = fitz.open(pdf_path)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        for page in doc:
            pix = page.get_pixmap(dpi=300)
            img_data = pix.tobytes("jpeg")
            
            # Add slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image from bytes
            img_stream = BytesIO(img_data)
            slide.shapes.add_picture(
                img_stream,
                Inches(0),
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
        
        prs.save(output_path)
        return True
    except Exception as e:
        print(f"Error: {str(e)}")
        return False

@pdf_to_ppt_bp.route('pdf_to_ppt')
def index():
    return render_template('pdf_to_ppt.html')

@pdf_to_ppt_bp.route('pdf_to_ppt/convert', methods=['POST'])
def convert_pdf_to_pptx_route():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if file and allowed_file(file.filename):
        # Generate unique filenames
        original_filename = secure_filename(file.filename)
        unique_id = str(uuid.uuid4())
        upload_path = os.path.join(pdf_to_ppt_bp.output_folder, f"{unique_id}_{original_filename}")
        output_filename = f"{os.path.splitext(original_filename)[0]}_{unique_id}.pptx"
        output_path = os.path.join(pdf_to_ppt_bp.output_folder, output_filename)
        
        # Save uploaded file
        file.save(upload_path)
        
        try:
            # Convert PDF to PPTX
            success = convert_pdf_to_pptx(upload_path, output_path)
            
            if not success:
                return jsonify({'error': 'Conversion failed'}), 500
            
            # Return success response with download link
            return jsonify({
                'success': True,
                'download_url': f'pdf_to_ppt/download/{output_filename}',
                'filename': output_filename
            })
        except Exception as e:
            return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
        finally:
            # Clean up uploaded file
            if os.path.exists(upload_path):
                os.remove(upload_path)
    else:
        return jsonify({'error': 'Invalid file type. Only PDF files are allowed.'}), 400

@pdf_to_ppt_bp.route('pdf_to_ppt/download/<filename>')
def download_file(filename):
    return send_from_directory(pdf_to_ppt_bp.output_folder, filename, as_attachment=True)

@pdf_to_ppt_bp.route('/cleanup', methods=['POST'])
def cleanup():
    try:
        # Clean up output files older than 1 hour
        now = datetime.now()
        for filename in os.listdir(pdf_to_ppt_bp.output_folder):
            file_path = os.path.join(pdf_to_ppt_bp.output_folder, filename)
            file_time = datetime.fromtimestamp(os.path.getmtime(file_path))
            if (now - file_time).total_seconds() > 3600:  # 1 hour
                os.remove(file_path)
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

